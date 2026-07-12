"""
build_presentation.py
Generates web-modernize-presentation.pptx — a 10-slide leadership deck.
Run: python build_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Slide geometry (16:9 widescreen) ────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Brand palette ───────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x2A)
COBALT  = RGBColor(0x1A, 0x56, 0xDB)
TEAL    = RGBColor(0x00, 0xC2, 0xCB)
WHITE   = RGBColor(0xFA, 0xFB, 0xFC)
SLATE   = RGBColor(0x6B, 0x72, 0x80)
INK     = RGBColor(0x1F, 0x29, 0x37)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
RED     = RGBColor(0xEF, 0x44, 0x44)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
LIGHT_B = RGBColor(0xEB, 0xF2, 0xFF)
LIGHT_G = RGBColor(0xF4, 0xF6, 0xF8)
BAND_BG = RGBColor(0xE6, 0xEC, 0xF2)

FONT_BODY        = "Calibri"
FONT_MONO        = "Consolas"
FONT_HEADER      = "Segoe UI Semibold"   # slide titles, section banners, table headers
FONT_HEADER_LIGHT = "Segoe UI Light"     # hero wordmarks on title/closing slides

FOOTER_TEXT = "web-modernize v0.18.0   ·   Balaji Harikrishnan   ·   Cognizant   ·   July 2026"


# ── Low-level helpers ───────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w if line_w else 0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, l, t, w, h, text, *, size=14, bold=False, italic=False,
             color=INK, font=FONT_BODY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def add_arrow_right(slide, l, t, w, h, fill=COBALT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ── Page chrome (header band + footer) ──────────────────────────────────────

def page_chrome(slide, title, subtitle=None, page_num=None):
    add_rect(slide, 0, 0, 13.333, 7.5, fill=LIGHT_G)
    add_rect(slide, 0, 0, 13.333, 1.05, fill=NAVY)
    # Finer left accent strip — refined feel
    add_rect(slide, 0, 0, 0.06, 1.05, fill=TEAL)
    # Title uses Segoe UI Semibold; the weight is already built into the face, no bold flag needed.
    add_text(slide, 0.30, 0.10, 12.8, 0.55, title,
             size=26, color=WHITE, font=FONT_HEADER,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, 0.30, 0.62, 12.8, 0.36, subtitle,
                 size=12, italic=True, color=TEAL)
    add_rect(slide, 0, 7.22, 13.333, 0.28, fill=BAND_BG)
    add_text(slide, 0.30, 7.24, 10.8, 0.24, FOOTER_TEXT,
             size=9, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
    if page_num is not None:
        add_text(slide, 12.30, 7.24, 0.85, 0.24,
                 f"{page_num:02d}  /  13",
                 size=9, color=SLATE, font=FONT_HEADER,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


# ── Table helper ────────────────────────────────────────────────────────────

def add_table(slide, l, t, w, h, headers, rows, *,
              col_ratios, header_size=12, body_size=11,
              mono_col0=False, header_fill=COBALT,
              alt_fill=LIGHT_B, row_height=None):
    ncols = len(headers)
    nrows = len(rows) + 1

    table_shape = slide.shapes.add_table(
        nrows, ncols, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tbl = table_shape.table

    total_emu = Inches(w)
    for i, ratio in enumerate(col_ratios):
        tbl.columns[i].width = int(total_emu * ratio)

    if row_height is not None:
        for i in range(nrows):
            tbl.rows[i].height = Inches(row_height)

    def style(cell, text, *, bold=False, fill=None, fg=INK,
              font=FONT_BODY, size=body_size, align=PP_ALIGN.LEFT,
              italic=False, anchor=MSO_ANCHOR.MIDDLE):
        cell.text = ""
        if fill is not None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
        cell.vertical_anchor = anchor
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = fg

    for ci, hdr in enumerate(headers):
        style(tbl.cell(0, ci), hdr, fill=header_fill,
              fg=WHITE, size=header_size, font=FONT_HEADER)

    for ri, row in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else alt_fill
        for ci, val in enumerate(row):
            font = FONT_MONO if (mono_col0 and ci == 0) else FONT_BODY
            color = COBALT if (mono_col0 and ci == 0) else INK
            bold = mono_col0 and ci == 0
            style(tbl.cell(ri + 1, ci), str(val),
                  fill=bg, fg=color, font=font, size=body_size,
                  bold=bold)

    return tbl


# ── Card grid helper (used by Need + Planning slides) ───────────────────────

def card_grid_3x2(slide, top_y, cards, accent_color):
    """Render 6 cards in a 3x2 grid. cards = [(title, body), ...]."""
    card_w = 4.04
    card_h = 2.45
    gap_x = 0.20
    gap_y = 0.20
    left_margin = 0.40

    for i, (title, body) in enumerate(cards):
        row = i // 3
        col = i % 3
        x = left_margin + col * (card_w + gap_x)
        y = top_y + row * (card_h + gap_y)
        # Card
        add_rect(slide, x, y, card_w, card_h, fill=WHITE,
                 line=accent_color, line_w=0.75)
        # Left accent strip
        add_rect(slide, x, y, 0.08, card_h, fill=accent_color)
        # Title
        add_text(slide, x + 0.22, y + 0.18, card_w - 0.34, 0.50,
                 title, size=14, bold=True, color=accent_color)
        # Body
        add_text(slide, x + 0.22, y + 0.70, card_w - 0.34, card_h - 0.85,
                 body, size=12, color=INK)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 13.333, 7.5, fill=NAVY)
    add_rect(slide, 0, 0, 13.333, 0.07, fill=TEAL)
    add_rect(slide, 0, 7.43, 13.333, 0.07, fill=TEAL)
    add_rect(slide, 0, 0, 0.14, 7.5, fill=TEAL)

    # Hero wordmark — Segoe UI Light at large size reads as elegant, not chunky
    add_text(slide, 0.5, 2.0, 12.333, 1.5, "web-modernize",
             size=72, color=WHITE, font=FONT_HEADER_LIGHT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Slightly wider teal rule, centred
    add_rect(slide, 4.42, 3.55, 4.50, 0.04, fill=TEAL)

    add_text(slide, 0.5, 3.75, 12.333, 0.55,
             "AI-Guided Legacy Web Application Migration",
             size=22, color=WHITE, font=FONT_HEADER, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 4.35, 12.333, 0.45,
             "A Claude Code Plugin for Enterprise Teams",
             size=15, italic=True, color=TEAL, align=PP_ALIGN.CENTER)

    add_text(slide, 0.5, 6.55, 12.333, 0.35,
             "Balaji Harikrishnan   ·   Cognizant   ·   July 2026   ·   v0.18.0",
             size=12, color=SLATE, font=FONT_HEADER, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 2 — What is web-modernize?  (plain English + technical)
# ══════════════════════════════════════════════════════════════════════════════

def slide_02_intro(prs):
    slide = blank_slide(prs)
    page_chrome(
        slide,
        "What is  web-modernize?",
        "A Claude Code plugin that gives software teams a structured workflow for modernizing legacy web applications",
        page_num=2,
    )

    # ─── IN PLAIN ENGLISH section ─────────────────────────────────
    add_rect(slide, 0.40, 1.20, 12.533, 0.40, fill=TEAL)
    add_text(slide, 0.55, 1.22, 12.333, 0.36,
             "OVERVIEW   —   what the plugin does",
             size=13, color=WHITE, font=FONT_HEADER,
             anchor=MSO_ANCHOR.MIDDLE)

    # Description card
    plain_text = (
        "web-modernize is a Claude Code plugin that developers install from "
        "the marketplace onto their own machine. It is used to modernize "
        "old web applications: the plugin reads a legacy application built "
        "in older technology (like ASP.NET, Java JSP, or AngularJS), then "
        "helps the team rewrite it  —  one page or feature at a time  —  "
        "into a modern framework (like React, Next.js, or Angular). "
        "Developers stay in control of every decision. The plugin handles "
        "the time-consuming translation work."
    )
    add_rect(slide, 0.40, 1.70, 12.533, 1.20, fill=WHITE,
             line=COBALT, line_w=0.75)
    add_text(slide, 0.60, 1.75, 12.133, 1.10, plain_text,
             size=13, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # Outcome chips
    chips = [
        "Migrates feature by feature, not all at once",
        "Visible to leadership at every step",
        "Team-scale workflow, not individual prompting",
    ]
    chip_w = 4.00
    chip_gap = 0.20
    total = chip_w * 3 + chip_gap * 2
    chip_start = (13.333 - total) / 2
    for i, txt in enumerate(chips):
        x = chip_start + i * (chip_w + chip_gap)
        add_rect(slide, x, 3.00, chip_w, 0.50, fill=COBALT)
        add_text(slide, x + 0.10, 3.02, chip_w - 0.20, 0.46,
                 txt, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ─── UNDER THE HOOD section ───────────────────────────────────
    add_rect(slide, 0.40, 3.65, 12.533, 0.40, fill=NAVY)
    add_text(slide, 0.55, 3.67, 12.333, 0.36,
             "TECHNICAL OVERVIEW   —   components and architecture",
             size=13, color=WHITE, font=FONT_HEADER,
             anchor=MSO_ANCHOR.MIDDLE)

    # 2 x 2 grid of technical cards
    tech_cards = [
        ("Plugin Host",
         "Runs inside Claude Code (Anthropic's official AI coding tool). "
         "Installs via the standard  /plugin install  command  —  no "
         "separate infrastructure, no servers."),
        ("What It Adds",
         "19 skills exposed as slash commands  ·  5 AI agents (analyzer, "
         "migrator, cross-cutting-migrator, parity-reviewer, migration-critic)  ·  "
         "2 automation hooks  ·  four workflows/ "
         "orchestration scripts  ·  9 schema-validated templates copied into the team's repo."),
        ("State and Concurrency",
         "All migration state lives in git as JSON. Per-unit file split "
         "(schema v3) lets 20+ developers migrate different units in "
         "parallel  —  zero merge conflicts by design."),
        ("Stacks Supported",
         "17 source stacks (ASP.NET WebForms/MVC/Core, Java JSP/Struts/"
         "Spring, AngularJS, Vue 2, jQuery, PHP, ColdFusion, Rails, "
         "Django, WordPress, ExtJS, classic ASP)  →  14 modern targets "
         "(React/Next/Vue 3/Angular/SvelteKit/Astro/Nuxt/Remix + "
         ".NET / Spring Boot / Nest / FastAPI / Express / Hono). Any "
         "other stack drops into a 3-question follow-up; migration "
         "continues."),
    ]
    card_w = 6.16
    card_h = 1.18
    gap_x = 0.20
    gap_y = 0.10
    base_x = 0.40
    base_y = 4.15

    for i, (title, body) in enumerate(tech_cards):
        row = i // 2
        col = i % 2
        x = base_x + col * (card_w + gap_x)
        y = base_y + row * (card_h + gap_y)
        # Card
        add_rect(slide, x, y, card_w, card_h, fill=WHITE,
                 line=NAVY, line_w=0.75)
        # Navy title strip
        add_rect(slide, x, y, card_w, 0.32, fill=NAVY)
        add_text(slide, x + 0.15, y + 0.02, card_w - 0.30, 0.30,
                 title, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        # Body
        add_text(slide, x + 0.15, y + 0.38, card_w - 0.30, card_h - 0.45,
                 body, size=11, color=INK,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # Takeaway strip
    add_rect(slide, 0.40, 6.78, 12.533, 0.36, fill=COBALT)
    add_text(slide, 0.50, 6.80, 12.333, 0.32,
             "Installed in minutes from the marketplace.  "
             "The team has a working, end-to-end migration workflow the same day.",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 3 — Why We Need It
# ══════════════════════════════════════════════════════════════════════════════

def slide_03_need(prs):
    slide = blank_slide(prs)
    page_chrome(
        slide,
        "Why We Need This Plugin",
        "Ad-hoc AI prompting and one-off tooling don't scale to enterprise modernization — six things the plugin adds",
        page_num=5,
    )

    advantages = [
        ("Enforced consistency across the team",
         "One workflow, one set of conventions, one quality bar. Every "
         "developer produces the same shape of output instead of inventing "
         "their own prompt style."),
        ("Durable progress ledger in git",
         "Who did what, what is done, what is left — all tracked. Direct "
         "prompting forgets everything when the session ends; the plugin "
         "remembers across days and developers."),
        ("Project context captured once, used everywhere",
         "Legacy stack, target stack, conventions, and constraints live in "
         "migration.md. No re-explaining the project on every prompt to "
         "every developer."),
        ("Verification gate that cannot be skipped",
         "Every unit must pass lint, type-check, and tests before it counts "
         "as done. Quality is enforced by the workflow, not remembered "
         "by the developer."),
        ("Team-aware coordination at scale",
         "Backlog, real-time kanban, per-unit git layout, deterministic "
         "merge rules. 20+ developers can migrate in parallel without "
         "colliding or duplicating work."),
        ("Built-in rollback and audit trail",
         "Every attempt, failure, and retry is recorded. One command "
         "(/rollback) reverts a unit cleanly. Direct prompting leaves you "
         "stranded when something breaks in production."),
    ]

    card_grid_3x2(slide, top_y=1.30, cards=advantages, accent_color=COBALT)

    # Closing pitch strip
    add_rect(slide, 0.40, 6.75, 12.533, 0.40, fill=NAVY)
    add_text(slide, 0.50, 6.78, 12.333, 0.34,
             "The plugin turns  'ask Claude to migrate this'  into a repeatable, "
             "auditable, team-scale engineering process.",
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 4 — From Legacy Codebase to Migration Backlog
# ══════════════════════════════════════════════════════════════════════════════

def slide_04_analyze(prs):
    slide = blank_slide(prs)
    page_chrome(
        slide,
        "From Legacy Codebase to Migration Backlog",
        "Two plugin commands  —  /analyze  and  /plan  —  produce a sized, prioritised backlog",
        page_num=8,
    )

    # Three-phase flow at the top
    box_w, box_h = 3.70, 1.20
    arrow_w, arrow_h = 0.80, 0.50
    flow_y = 1.30
    xs = [0.32, 4.82, 9.32]  # box left positions

    phase_titles = [
        ("Legacy Codebase", "Untouched.\nRead-only inspection."),
        ("Pre-filled  migration.md", "Auto-detected stack +\nteam-edited targets."),
        ("Migration Backlog", "Sized, prioritised,\nready for sprint planning."),
    ]
    for i, (title, sub) in enumerate(phase_titles):
        x = xs[i]
        # Box
        fill_c = NAVY if i == 1 else COBALT
        add_rect(slide, x, flow_y, box_w, box_h, fill=fill_c)
        add_text(slide, x + 0.10, flow_y + 0.15, box_w - 0.20, 0.40,
                 title, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + 0.10, flow_y + 0.55, box_w - 0.20, 0.60,
                 sub, size=11, italic=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Arrows between boxes
    arrow_labels = ["/analyze", "/plan"]
    arrow_xs = [4.05, 8.55]
    for i, (ax, label) in enumerate(zip(arrow_xs, arrow_labels)):
        add_arrow_right(slide, ax, flow_y + 0.35, arrow_w, arrow_h, fill=TEAL)
        add_text(slide, ax - 0.05, flow_y + 0.93, arrow_w + 0.10, 0.30,
                 label, size=12, bold=True, color=TEAL,
                 font=FONT_MONO, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Three explanatory cards underneath
    cards_y = 2.85
    cards_h = 3.85
    card_w = 4.04
    gap = 0.20
    left = 0.40

    phase_cards = [
        ("Step 1  —  /analyze (automatic)", COBALT, [
            "Reads every legacy file (never modifies them).",
            "Detects framework, version, and build tooling.",
            "Identifies entry points  —  exhaustive loop-until-dry discovery on large apps.",
            "Pre-fills migration.md §2  (Source Stack)  for you.",
            "Writes analysis.json — a structured inventory you can audit.",
        ]),
        ("Step 2  —  Edit migration.md (team)", TEAL, [
            "Choose target UI framework  (React, Next, Vue, Angular, Svelte).",
            "Choose target API framework  —  or 'none' for UI-only migrations.",
            "Pick migration strategy  (strangler-fig, big-bang, module-by-module).",
            "Declare current and target auth provider.",
            "Set acceptance criteria  —  these drive /verify's pass/fail bar.",
        ]),
        ("Step 3  —  /plan (automatic)", GREEN, [
            "Validates that migration.md is complete and consistent.",
            "Generates plan.md  —  a reviewable plan with a dependency graph.",
            "Seeds the unit backlog  —  one item per page or feature.",
            "Re-runnable: edit migration.md and re-run /plan anytime.",
            "History is preserved across re-plans  (by unit id).",
        ]),
    ]

    for i, (title, color, bullets) in enumerate(phase_cards):
        x = left + i * (card_w + gap)
        # Card border
        add_rect(slide, x, cards_y, card_w, cards_h, fill=WHITE,
                 line=color, line_w=0.75)
        # Title strip
        add_rect(slide, x, cards_y, card_w, 0.55, fill=color)
        add_text(slide, x + 0.10, cards_y + 0.05, card_w - 0.20, 0.45,
                 title, size=13, color=WHITE, font=FONT_HEADER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Bullets
        bullet_y = cards_y + 0.70
        for bullet in bullets:
            add_text(slide, x + 0.20, bullet_y, 0.20, 0.30,
                     "▸", size=12, bold=True, color=color)
            add_text(slide, x + 0.42, bullet_y, card_w - 0.55, 0.55,
                     bullet, size=11, color=INK)
            bullet_y += 0.60

    # Bottom takeaway strip
    add_rect(slide, 0.40, 6.78, 12.533, 0.36, fill=NAVY)
    add_text(slide, 0.50, 6.80, 12.333, 0.32,
             "All artifacts  (migration.md, analysis.json, plan.md, units/*.json)  "
             "live in git  —  auditable, versioned, never lost to a laptop.",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 5 — From Backlog to Migrated Units
# ══════════════════════════════════════════════════════════════════════════════

def slide_05_execution(prs):
    slide = blank_slide(prs)
    page_chrome(
        slide,
        "From Backlog to Migrated Units",
        "Set up once with  /scaffold  and  /foundation.  Then the team migrates in parallel with  /next  and  /migrate.",
        page_num=9,
    )

    # ─── SETUP section banner ─────────────────────────────────────
    add_rect(slide, 0.40, 1.30, 12.533, 0.42, fill=SLATE)
    add_text(slide, 0.55, 1.32, 12.333, 0.38,
             "SETUP   —   one developer runs   /scaffold   then   /foundation",
             size=13, color=WHITE, font=FONT_HEADER,
             anchor=MSO_ANCHOR.MIDDLE)

    # 3-box linear flow (mirrors slide 3 geometry)
    box_w, box_h = 3.70, 1.10
    arrow_w, arrow_h = 0.80, 0.40
    flow_y = 1.92
    xs = [0.32, 4.82, 9.32]

    setup_states = [
        ("Backlog Ready",   "From  /plan.\nUnits sized S / M / L / XL."),
        ("Target Scaffold", "Skeleton + toolchain preflight.\nLegacy assets copied in."),
        ("Foundation Ready", "Auth + cross-cutting slice\nmigrated first."),
    ]
    for i, (title, sub) in enumerate(setup_states):
        x = xs[i]
        # Middle box highlighted to draw the eye to the in-progress state
        fill_c = NAVY if i == 1 else COBALT
        add_rect(slide, x, flow_y, box_w, box_h, fill=fill_c)
        add_text(slide, x + 0.10, flow_y + 0.10, box_w - 0.20, 0.40,
                 title, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + 0.10, flow_y + 0.50, box_w - 0.20, 0.55,
                 sub, size=10, italic=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Arrows between setup boxes
    arrow_labels = ["/scaffold", "/foundation"]
    arrow_xs = [4.05, 8.55]
    for ax, label in zip(arrow_xs, arrow_labels):
        add_arrow_right(slide, ax, flow_y + 0.35, arrow_w, arrow_h, fill=TEAL)
        add_text(slide, ax - 0.05, flow_y + 0.78, arrow_w + 0.10, 0.30,
                 label, size=11, bold=True, color=TEAL,
                 font=FONT_MONO, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # ─── PARALLEL MIGRATION section banner ────────────────────────
    add_rect(slide, 0.40, 3.40, 12.533, 0.42, fill=TEAL)
    add_text(slide, 0.55, 3.42, 12.333, 0.38,
             "PARALLEL MIGRATION   —   each developer runs   /next   or   /migrate   then   /verify   on a different unit",
             size=13, color=WHITE, font=FONT_HEADER,
             anchor=MSO_ANCHOR.MIDDLE)

    # 5 dev cards in a row — each developer migrates a different unit at the same time
    dev_cards = [
        ("Dev A", "HomePage",         "/migrate  →  /verify"),
        ("Dev B", "PaymentProcessor", "/migrate  →  /verify"),
        ("Dev C", "OrderList",        "/migrate  →  /verify"),
        ("Dev D", "Dashboard",        "/migrate  →  /verify"),
        ("Dev E", "ProfileEditor",    "/migrate  →  /verify"),
    ]
    # Single corporate accent — distinction between developers comes from labels
    # and unit names, not rainbow colors.
    card_w = 2.43
    card_h = 2.40
    gap = 0.13
    cards_y = 4.00
    left = 0.40

    for i, (dev, unit, cmd) in enumerate(dev_cards):
        x = left + i * (card_w + gap)
        color = COBALT
        # Card body
        add_rect(slide, x, cards_y, card_w, card_h, fill=WHITE,
                 line=color, line_w=0.75)
        # Top color strip
        add_rect(slide, x, cards_y, card_w, 0.45, fill=color)
        add_text(slide, x + 0.10, cards_y + 0.04, card_w - 0.20, 0.40,
                 dev, size=13, color=WHITE, font=FONT_HEADER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # "migrating" label
        add_text(slide, x + 0.10, cards_y + 0.60, card_w - 0.20, 0.30,
                 "migrating", size=10, italic=True, color=SLATE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Unit name
        add_text(slide, x + 0.10, cards_y + 0.95, card_w - 0.20, 0.55,
                 unit, size=13, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Command at bottom
        add_rect(slide, x + 0.12, cards_y + 1.70, card_w - 0.24, 0.50,
                 fill=LIGHT_B, line=color, line_w=0.5)
        add_text(slide, x + 0.12, cards_y + 1.72, card_w - 0.24, 0.46,
                 cmd, size=10, bold=True, color=color,
                 font=FONT_MONO, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # "Same time" indicator below cards
    add_text(slide, 0.40, 6.48, 12.533, 0.26,
             "/verify enforces lint + type-check + tests + behavioural-parity + security checks (plus an advisory code-quality review) before a unit is marked done.  Per-unit git files mean zero merge conflicts.",
             size=11, italic=True, color=SLATE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ─── Takeaway strip ───────────────────────────────────────────
    add_rect(slide, 0.40, 6.78, 12.533, 0.36, fill=NAVY)
    add_text(slide, 0.50, 6.80, 12.333, 0.32,
             "Setup is sequential and short.  Migration is parallel and scales with team size.",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 8 — Risk, Safety, and Reversibility
# ══════════════════════════════════════════════════════════════════════════════

def slide_08_risk_safety(prs):
    slide = blank_slide(prs)
    page_chrome(
        slide,
        "Risk, Safety, and Reversibility",
        "What if the AI gets it wrong?  —  every change is reviewable, reversible, and auditable",
        page_num=6,
    )

    cards = [
        ("Approved going in, reversible coming out",
         "Nothing risky is written blind: by default /next presents each "
         "unit's plan — target files, approach, tests — and waits for your "
         "approval before writing a line (opt out per migration via "
         "review_mode, or per unit via --no-plan); /foundation always gates "
         "the auth + cross-cutting design. And nothing is one-way: every "
         "action is recorded in git, one command  (/rollback --unit <id>)  "
         "reverts a unit cleanly (refusing shared files unless forced), and "
         "destructive commands  (/abandon --hard)  need a two-step confirm. "
         "Per-unit notes capture the why alongside the what."),
        ("Nothing ships unverified",
         "/verify is a hard gate — lint, type-check, and tests must pass "
         "before a unit transitions from migrated to verified. A read-only "
         "parity-reviewer then compares migrated vs legacy — validation, "
         "output shape, errors, UI states, and security (dropped "
         "authorization, injection, secret leakage); high-severity "
         "differences block until acknowledged via /parity-check. An "
         "advisory migration-critic flags non-idiomatic code and "
         "static-performance regressions (N+1 queries, unbounded fetches, "
         "bundle bloat), and an opt-in dynamic tier (API replay + Playwright "
         "E2E) adds runtime checks. Discovered credentials are masked in "
         "every shared artifact. /retry --with-prompt re-attempts with "
         "human guidance."),
        ("Nothing is hidden",
         "All state lives in git as JSON — auditable, diffable, "
         "code-reviewable. /status shows real-time who's working on "
         "what; stale sessions surface automatically. Per-unit file "
         "split means each developer's work appears as its own commit; "
         "nothing hides in a merge."),
    ]

    # 3-column card grid (re-uses card_grid_3x2's left-accent visual pattern,
    # but 3-across instead of 3x2)
    card_w = 4.04
    card_h = 4.85
    gap_x = 0.20
    left_margin = 0.40
    top_y = 1.30
    accent_colors = [COBALT, TEAL, GREEN]

    for i, (title, body) in enumerate(cards):
        x = left_margin + i * (card_w + gap_x)
        color = accent_colors[i]
        # Card body
        add_rect(slide, x, top_y, card_w, card_h, fill=WHITE,
                 line=color, line_w=0.75)
        # Left accent strip
        add_rect(slide, x, top_y, 0.10, card_h, fill=color)
        # Title
        add_text(slide, x + 0.25, top_y + 0.20, card_w - 0.40, 0.55,
                 title, size=15, bold=True, color=color)
        # Body
        add_text(slide, x + 0.25, top_y + 0.85, card_w - 0.40, card_h - 1.05,
                 body, size=12, color=INK)

    # Bottom takeaway band
    add_rect(slide, 0.40, 6.40, 12.533, 0.75, fill=NAVY)
    add_text(slide, 0.50, 6.42, 12.333, 0.36,
             "Every AI-generated change is reviewable, reversible, and auditable",
             size=15, bold=True, color=WHITE, font=FONT_HEADER,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 0.50, 6.78, 12.333, 0.34,
             "Not a one-way bet on the model getting it right the first time.",
             size=12, italic=True, color=TEAL,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 9 — Natural Language + Pluggable Framework Library
# ══════════════════════════════════════════════════════════════════════════════

def slide_09_extensibility(prs):
    slide = blank_slide(prs)
    page_chrome(
        slide,
        "Natural Language  &  Pluggable Framework Library",
        "Two improvements in v0.10.0 that make the plugin easier to use and easier to extend",
        page_num=3,
    )

    # ─── LEFT HALF: Natural-Language Routing ──────────────────────
    add_rect(slide, 0.40, 1.30, 6.16, 0.42, fill=COBALT)
    add_text(slide, 0.55, 1.32, 6.00, 0.38,
             "TALK TO IT NATURALLY   —   no slash commands required",
             size=12, color=WHITE, font=FONT_HEADER,
             anchor=MSO_ANCHOR.MIDDLE)

    # Description
    add_rect(slide, 0.40, 1.78, 6.16, 0.95, fill=WHITE,
             line=COBALT, line_w=0.75)
    add_text(slide, 0.55, 1.83, 5.86, 0.87,
             "Every skill's description now packs intent phrases + a "
             "lifecycle anchor. Claude's native skill matcher routes "
             "plain-English requests to the right command — developers "
             "stop memorising 19 slash commands.",
             size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # Examples table — "you type → skill fires"
    nl_examples = [
        ("“what's next”",            "/next"),
        ("“let's plan it”",          "/plan"),
        ("“where are we”",           "/status"),
        ("“migrate the login page”", "/migrate"),
        ("“stuck lock”",             "/unlock"),
        ("“start over”",             "/abandon"),
    ]
    ex_y = 2.85
    ex_h = 0.35
    # Header
    add_rect(slide, 0.40, ex_y, 3.20, ex_h, fill=NAVY)
    add_text(slide, 0.50, ex_y + 0.02, 3.00, ex_h - 0.04,
             "You type", size=11, color=WHITE, font=FONT_HEADER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 3.60, ex_y, 2.96, ex_h, fill=NAVY)
    add_text(slide, 3.70, ex_y + 0.02, 2.76, ex_h - 0.04,
             "Skill that fires", size=11, color=WHITE, font=FONT_HEADER,
             anchor=MSO_ANCHOR.MIDDLE)
    for i, (utterance, skill) in enumerate(nl_examples):
        y = ex_y + ex_h + i * ex_h
        bg = WHITE if i % 2 == 0 else LIGHT_B
        add_rect(slide, 0.40, y, 3.20, ex_h, fill=bg)
        add_text(slide, 0.50, y + 0.02, 3.00, ex_h - 0.04,
                 utterance, size=11, color=INK, italic=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, 3.60, y, 2.96, ex_h, fill=bg)
        add_text(slide, 3.70, y + 0.02, 2.76, ex_h - 0.04,
                 skill, size=11, color=COBALT, bold=True,
                 font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)

    # ─── RIGHT HALF: Pluggable Framework Library ──────────────────
    add_rect(slide, 6.77, 1.30, 6.16, 0.42, fill=TEAL)
    add_text(slide, 6.92, 1.32, 6.00, 0.38,
             "PLUGGABLE FRAMEWORKS   —   drop a file, get a new stack",
             size=12, color=WHITE, font=FONT_HEADER,
             anchor=MSO_ANCHOR.MIDDLE)

    # Description
    add_rect(slide, 6.77, 1.78, 6.16, 0.95, fill=WHITE,
             line=TEAL, line_w=0.75)
    add_text(slide, 6.92, 1.83, 5.86, 0.87,
             "31 frameworks/<name>.md files — one per supported "
             "stack — hold the recipes /scaffold, /foundation, and "
             "legacy-analyzer all read on demand. Adding a new "
             "framework is a one-file drop-in.",
             size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 4 capability badges — what comes from the framework library
    badges = [
        ("17",  "Legacy source stacks auto-detected"),
        ("14",  "Modern targets with full scaffold recipes"),
        ("1",   "File to add a brand-new framework"),
        ("0",   "Plugin code changes needed for the new framework"),
    ]
    badge_w = 2.95
    badge_h = 0.85
    badge_gap = 0.10
    badge_start_x = 6.77
    badge_start_y = 2.85
    for i, (number, label) in enumerate(badges):
        row = i // 2
        col = i % 2
        x = badge_start_x + col * (badge_w + badge_gap)
        y = badge_start_y + row * (badge_h + badge_gap)
        add_rect(slide, x, y, badge_w, badge_h, fill=WHITE,
                 line=TEAL, line_w=0.75)
        # Big number on the left
        add_rect(slide, x, y, 0.80, badge_h, fill=TEAL)
        add_text(slide, x, y, 0.80, badge_h,
                 number, size=28, color=WHITE, font=FONT_HEADER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Label on the right
        add_text(slide, x + 0.90, y + 0.05, badge_w - 1.00, badge_h - 0.10,
                 label, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ─── UNKNOWN-TECH PATH — runs full width across the lower section ─
    add_rect(slide, 0.40, 4.85, 12.53, 0.42, fill=NAVY)
    add_text(slide, 0.55, 4.87, 12.30, 0.38,
             "UNKNOWN-TECH PATH   —   when the framework isn't in the library yet",
             size=12, color=WHITE, font=FONT_HEADER,
             anchor=MSO_ANCHOR.MIDDLE)

    # Three-step flow showing what happens with an unknown framework
    box_w, box_h = 3.95, 1.40
    arrow_w, arrow_h = 0.40, 0.40
    flow_y = 5.36
    box_xs = [0.40, 4.69, 8.98]
    flow_titles = [
        ("Source unknown?",
         "/analyze shows the raw evidence (file extensions, "
         "library refs, build files) and asks you to name the stack. "
         "Migration continues with the user-supplied name."),
        ("Target unknown?",
         "/scaffold asks 3 questions: scaffold command, test "
         "framework, verify commands. Answers persist in "
         "verify.config.json so retries don't re-ask."),
        ("Auth unknown?",
         "/foundation skips the prebuilt password-hashing template, "
         "defers to permanent-gotchas + OWASP guidance. Auth "
         "migration completes; the team supplies code."),
    ]
    for i, (title, body) in enumerate(flow_titles):
        x = box_xs[i]
        add_rect(slide, x, flow_y, box_w, box_h, fill=WHITE,
                 line=NAVY, line_w=0.75)
        add_rect(slide, x, flow_y, box_w, 0.32, fill=NAVY)
        add_text(slide, x + 0.10, flow_y + 0.02, box_w - 0.20, 0.30,
                 title, size=12, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + 0.12, flow_y + 0.40, box_w - 0.24, box_h - 0.48,
                 body, size=10, color=INK, anchor=MSO_ANCHOR.TOP)

    # ─── Takeaway strip ───────────────────────────────────────────
    add_rect(slide, 0.40, 6.92, 12.533, 0.24, fill=COBALT)
    add_text(slide, 0.50, 6.93, 12.333, 0.22,
             "Friendlier to onboarding.  Quieter when the right answer is obvious.  Never stuck on an unsupported stack.",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 10 — Built for Sprint & PI Planning
# ══════════════════════════════════════════════════════════════════════════════

def slide_10_planning(prs):
    slide = blank_slide(prs)
    page_chrome(
        slide,
        "Built for Sprint and PI Planning",
        "The plugin speaks backlog, dependencies, sprints, and burndown natively",
        page_num=7,
    )

    pairings = [
        ("/plan  =  Sprint or PI planning input",
         "Auto-generates the full migration backlog. One estimable item per page "
         "or feature — ready for sizing, prioritisation, and assignment."),
        ("Each unit  =  one estimable backlog item",
         "Sized S / M / L / XL by file count and complexity. Discrete and "
         "independently testable, with a clear definition of done "
         "(/verify must pass). XL units become epics — teams decompose "
         "them into sprint-sized work during PI planning."),
        ("Dependencies built in",
         "/plan declares which units depend on which — a natural sequencing "
         "input for sprint commitment and PI iteration goals."),
        ("Standup-friendly coordination",
         "Teams pick assignments at standup, just like any sprint. Behind "
         "the scenes the plugin stores each unit's progress in its own "
         "file — so two developers on different units never collide in git."),
        ("/status  =  real-time team kanban",
         "Read-only view showing pending, in-flight, blocked, and done units "
         "across the whole team. Refreshes on every git pull."),
        ("/report  =  burndown and ETA",
         "Run /report whenever you need a burndown chart, velocity trend, "
         "or completion ETA — perfect for sprint reviews and PI demos. "
         "Numbers come from the team's actual git history, not a "
         "spreadsheet someone has to maintain."),
    ]

    card_grid_3x2(slide, top_y=1.30, cards=pairings, accent_color=COBALT)

    # Bottom Scaled Agile banner
    add_rect(slide, 0.40, 6.75, 12.533, 0.40, fill=NAVY)
    add_text(slide, 0.50, 6.78, 12.333, 0.34,
             "Designed for Scaled Agile  —  multiple feature teams in one ART "
             "can migrate different units in the same PI without colliding.",
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 6 — Three Verification Tiers
# ══════════════════════════════════════════════════════════════════════════════

def slide_06_verification_tiers(prs):
    slide = blank_slide(prs)
    page_chrome(
        slide,
        "Three Verification Tiers",
        "/verify runs all three in order — each tier has a different role, depth, and gate behaviour",
        page_num=10,
    )

    # Section banner
    add_rect(slide, 0.40, 1.30, 12.533, 0.42, fill=NAVY)
    add_text(slide, 0.55, 1.32, 12.333, 0.38,
             "Every migrated unit passes through up to three verification layers  ·  "
             "Tier 1 blocks  ·  Tier 2 advises  ·  Tier 3 adds runtime confidence (opt-in)",
             size=12, color=WHITE, font=FONT_HEADER,
             anchor=MSO_ANCHOR.MIDDLE)

    card_w   = 4.04
    card_gap = 0.20
    card_top = 1.82
    card_h   = 4.82
    left     = 0.40

    tiers = [
        {
            "num": "1",
            "label": "Parity + Security",
            "tag": "GATING",
            "accent": RED,
            "card_fill": RGBColor(0xFD, 0xEB, 0xEB),
            "description": (
                "A read-only AI (parity-reviewer) compares migrated code vs the legacy "
                "original across validation logic, output shape, error handling, UI "
                "structure, and five security dimensions: dropped auth, injection, "
                "encoding, CSRF, and exposed secrets."
            ),
            "bullets": [
                "Behavioural diff: logic, responses, error shapes match legacy",
                "Security: auth, injection, encoding, CSRF, secret leakage",
                "Adversarial refute pass on every high-severity finding",
                "Unacknowledged highs BLOCK the unit from being verified",
            ],
            "commands": "/verify  ·  /parity-check",
        },
        {
            "num": "2",
            "label": "Quality + Static Perf",
            "tag": "ADVISORY",
            "accent": COBALT,
            "card_fill": LIGHT_B,
            "description": (
                "A read-only AI (migration-critic) reviews for two orthogonal concerns: "
                "idiomatic quality (legacy patterns leaking into modern code, e.g. jQuery "
                "patterns in React) and static-performance regressions detectable "
                "without running the app."
            ),
            "bullets": [
                "Idiomatic review: legacy paradigm leakage in modern code",
                "N+1 queries, unbounded data fetches, waterfall I/O",
                "Bundle bloat and blocking synchronous operations",
                "Findings reported to the developer — never blocks migration",
            ],
            "commands": "/verify  ·  /quality-check",
        },
        {
            "num": "3",
            "label": "Dynamic Testing",
            "tag": "OPT-IN",
            "accent": PURPLE,
            "card_fill": RGBColor(0xF3, 0xE8, 0xFF),
            "description": (
                "Runtime verification enabled per migration in migration.md §12 or with "
                "--dynamic. Adds two phases on top of the static tiers: recorded-baseline "
                "API replay against the new API, and browser-driven Playwright "
                "end-to-end tests."
            ),
            "bullets": [
                "API replay: records legacy responses, diffs new API outputs",
                "Playwright E2E: browser-driven suite vs the running new app",
                "Enable with /verify --dynamic or migration.md §12",
                "Run --capture-baseline before the first comparison",
            ],
            "commands": "/verify --dynamic  ·  --capture-baseline",
        },
    ]

    for i, tier in enumerate(tiers):
        x      = left + i * (card_w + card_gap)
        accent = tier["accent"]

        # Card body
        add_rect(slide, x, card_top, card_w, card_h, fill=tier["card_fill"],
                 line=accent, line_w=0.75)

        # Top accent strip
        add_rect(slide, x, card_top, card_w, 0.60, fill=accent)
        # Tier number (large)
        add_text(slide, x + 0.10, card_top + 0.02, 0.52, 0.56,
                 tier["num"], size=26, bold=True, color=WHITE,
                 font=FONT_HEADER, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Tier name
        add_text(slide, x + 0.72, card_top + 0.04, card_w - 0.84, 0.32,
                 tier["label"], size=13, bold=True, color=WHITE,
                 font=FONT_HEADER, anchor=MSO_ANCHOR.MIDDLE)
        # Tag badge row
        add_text(slide, x + 0.72, card_top + 0.36, card_w - 0.84, 0.20,
                 tier["tag"], size=9, bold=True, color=WHITE,
                 font=FONT_HEADER, anchor=MSO_ANCHOR.MIDDLE)

        # Description paragraph
        add_text(slide, x + 0.20, card_top + 0.70, card_w - 0.35, 1.48,
                 tier["description"], size=10.5, color=INK, anchor=MSO_ANCHOR.TOP)

        # Four bullet points
        bullet_y = card_top + 2.26
        for bullet in tier["bullets"]:
            add_text(slide, x + 0.20, bullet_y, 0.22, 0.40,
                     "▸", size=10, bold=True, color=accent)
            add_text(slide, x + 0.44, bullet_y, card_w - 0.58, 0.44,
                     bullet, size=10, color=INK)
            bullet_y += 0.46

        # Commands badge at bottom of card
        cmd_top = card_top + card_h - 0.52
        add_rect(slide, x + 0.12, cmd_top, card_w - 0.24, 0.42, fill=accent)
        add_text(slide, x + 0.12, cmd_top + 0.02, card_w - 0.24, 0.38,
                 tier["commands"], size=10, bold=True, color=WHITE,
                 font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom takeaway strip
    add_rect(slide, 0.40, 6.78, 12.533, 0.36, fill=NAVY)
    add_text(slide, 0.50, 6.80, 12.333, 0.32,
             "Tier 1 guards correctness and security.  Tier 2 guards code quality and performance.  "
             "Tier 3 adds runtime confidence.  Any combination is valid.",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 7 — What the Plugin Adds to Claude Code (before / after)
# ══════════════════════════════════════════════════════════════════════════════

def slide_07_advantages(prs):
    slide = blank_slide(prs)
    page_chrome(
        slide,
        "What the Plugin Adds to Claude Code",
        "Every developer already has Claude Code — here is what changes when you add web-modernize",
        page_num=4,
    )

    # 3-column before / after comparison
    headers = [
        "Category",
        "Without the plugin   (raw Claude Code)",
        "With   web-modernize",
    ]
    rows = [
        ("Setup ergonomics",
         "Developer hand-edits a config file: which target framework? which strategy? which auth? Easy to miss a field, no validation until later.",
         "/analyze walks the team through every required choice with stack-aware recommendations — pick from a list, the plugin writes the config."),
        ("Tool discovery",
         "Developer memorises slash commands or hunts through docs — 'what was the command to retry a failed unit again?'",
         "Type plain English — 'what's next', 'let's plan', 'stuck lock' — and the right skill auto-fires. No memorisation."),
        ("Defining the work",
         "Each developer decides what to migrate next based on chat or a shared doc — no sized backlog.",
         "/plan generates a sized, dependency-ordered backlog — every unit is estimable and assignable."),
        ("Project context",
         "Re-explain the legacy stack, target stack, and conventions on every prompt and every new session.",
         "Captured once in migration.md — every command reads it automatically, every developer stays aligned."),
        ("Team coordination",
         "20+ developers prompting in parallel produce merge conflicts and silently duplicated work.",
         "Per-unit git files mean Alice and Bob's work touch zero shared files — conflict-free by design."),
        ("Quality assurance",
         "Verification is whatever the developer remembers to run — lint, types, tests can be skipped.",
         "/verify enforces lint, type-check, tests, and a behavioural-parity check as a hard gate before a unit is marked done."),
        ("Failure recovery",
         "A failed migration leaves broken files behind — recovery is manual git surgery, no audit trail.",
         "/rollback reverts a unit cleanly in one command; /retry preserves every diagnostic for review."),
        ("Leadership reporting",
         "Progress lives in chat and spreadsheets — burndown is manual, ETA is a guess, risks are anecdotal.",
         "/report generates burndown, velocity, and ETA from git state — always accurate, no upkeep."),
    ]

    table_l, table_t = 0.40, 1.30
    table_w, table_h = 12.533, 5.40

    ncols = 3
    nrows = len(rows) + 1
    col_ratios = [0.18, 0.41, 0.41]
    col_widths = [table_w * r for r in col_ratios]
    col_xs = [table_l]
    for cw in col_widths[:-1]:
        col_xs.append(col_xs[-1] + cw)
    row_h = table_h / nrows

    # Header row
    for ci, hdr in enumerate(headers):
        # Header colors: category = navy, without = slate-ish, with = cobalt (winner)
        if ci == 0:
            fill = NAVY
        elif ci == 1:
            fill = SLATE
        else:
            fill = COBALT
        add_rect(slide, col_xs[ci], table_t, col_widths[ci], row_h, fill=fill)
        add_text(slide, col_xs[ci] + 0.10, table_t + 0.04,
                 col_widths[ci] - 0.20, row_h - 0.08, hdr,
                 size=13, color=WHITE, font=FONT_HEADER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Data rows
    for ri, row in enumerate(rows):
        y = table_t + row_h * (ri + 1)
        category, without_txt, with_txt = row
        alt = (ri % 2 == 1)

        # Category column
        cat_bg = WHITE if not alt else LIGHT_G
        add_rect(slide, col_xs[0], y, col_widths[0], row_h, fill=cat_bg)
        add_text(slide, col_xs[0] + 0.12, y + 0.05,
                 col_widths[0] - 0.24, row_h - 0.10, category,
                 size=12, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)

        # "Without" column — muted gray fill
        without_bg = LIGHT_G if not alt else BAND_BG
        add_rect(slide, col_xs[1], y, col_widths[1], row_h, fill=without_bg)
        add_text(slide, col_xs[1] + 0.15, y + 0.08,
                 col_widths[1] - 0.30, row_h - 0.16, without_txt,
                 size=11, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)

        # "With" column — cobalt-tinted, emphasised
        with_bg = LIGHT_B
        add_rect(slide, col_xs[2], y, col_widths[2], row_h,
                 fill=with_bg, line=COBALT, line_w=0.75)
        add_text(slide, col_xs[2] + 0.15, y + 0.08,
                 col_widths[2] - 0.30, row_h - 0.16, with_txt,
                 size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom takeaway strip
    add_rect(slide, 0.40, 6.78, 12.533, 0.36, fill=NAVY)
    add_text(slide, 0.50, 6.80, 12.333, 0.32,
             "Same Claude Code engine.  Structured workflow on top.  "
             "Every advantage above maps to a concrete plugin feature.",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 11 — The 19 Skills (Slash Commands)
# ══════════════════════════════════════════════════════════════════════════════

def slide_11_commands(prs):
    slide = blank_slide(prs)
    page_chrome(
        slide,
        "The 19 Skills  (Slash Commands)",
        "Each skill is exposed as a slash command — or auto-fires from plain English (see slide 3)",
        page_num=11,
    )

    headers = ["Command", "What It Does (Plain English)"]
    rows = [
        ("/init",      "Sets up the migration workspace inside the legacy repository."),
        ("/analyze",   "Auto-detects the legacy technology stack and inventories every component."),
        ("/plan",      "Validates the config and generates the full migration backlog — one item per page or feature."),
        ("/scaffold",  "Creates the blank modern application skeleton and copies legacy images, fonts, and assets."),
        ("/foundation","Establishes the foundational slice first — login/auth plus opted-in cross-cutting concerns (i18n, feature flags, error handling, telemetry, logging, data-access wiring) — via one consolidated design-approval gate, implemented in parallel."),
        ("/next",      "Picks the next available unit from the backlog and migrates it — presenting a plan for approval first (opt-out)."),
        ("/migrate",   "Migrates a specific unit by name — used when standup assigns work."),
        ("/verify",    "Runs lint, type-check, tests, a behavioural-parity + security gate, and advisory code-quality + static-performance review; an opt-in dynamic tier (API replay + Playwright E2E) adds runtime checks."),
        ("/parity-check", "Compares a migrated unit's behaviour against the legacy original — validation, output shape, errors, UI states, security — and lets the team acknowledge intentional differences."),
        ("/quality-check", "Advisory review of the migrated code's idiomaticity and static performance — flags legacy patterns and N+1/bundle/blocking regressions. Never blocks."),
        ("/integrate", "Assembles the migrated units into the composed app — central router + nav, whole-app smoke, orphaned-unit + cutover-coverage report, and the strangler traffic-split proxy. Runs any time; idempotent."),
        ("/retry",     "Re-attempts a failed unit, optionally with team-supplied corrective guidance."),
        ("/rollback",  "Reverts a single unit back to its original state in one git command — refusing by default when it owns shared files others depend on."),
        ("/sync",      "Merges teammates' parallel progress into the local copy with deterministic rules."),
        ("/report",    "Generates a stakeholder report — burndown, ETA, risk — in Markdown, JSON, or HTML."),
        ("/status",    "Shows a real-time dashboard of every unit, owner, and blocker (read-only)."),
        ("/unlock",    "Force-clears a stuck advisory lock left behind by a crashed session."),
        ("/abandon",   "Formally drops a unit or resets the workspace — requires two-step confirmation."),
    ]

    add_table(
        slide, 0.40, 1.26, 12.533, 5.95,
        headers, rows,
        col_ratios=[0.18, 0.82],
        header_size=13, body_size=10,
        mono_col0=True,
        row_height=0.305,
    )


# ══════════════════════════════════════════════════════════════════════════════
# Slide 12 — Agents · Hooks · Templates · Framework Library
# ══════════════════════════════════════════════════════════════════════════════

def slide_12_inventory(prs):
    slide = blank_slide(prs)
    page_chrome(
        slide,
        "Agents · Hooks · Templates · Framework Library",
        "Every other custom artifact built into the plugin — one English sentence each",
        page_num=12,
    )

    add_text(slide, 0.40, 1.20, 12.5, 0.30,
             "5 AI Agents   (specialised AIs that read your code, translate it, verify behaviour, quality & performance, and steer around known pitfalls)",
             size=14, color=COBALT, font=FONT_HEADER)
    add_table(
        slide, 0.40, 1.50, 12.533, 1.60,
        ["Agent", "What It Does"],
        [
            ("legacy-analyzer",
             "Read-only AI agent that inspects the legacy codebase and produces a structured stack inventory."),
            ("unit-migrator",
             "Core migration engine that translates one page or feature from legacy code to the modern target stack."),
            ("cross-cutting-migrator",
             "Establishes one cross-cutting concern (auth / i18n / flags / error-handling / telemetry / logging / data-wiring); fanned out in parallel by /foundation."),
            ("parity-reviewer",
             "Read-only AI that compares migrated vs legacy for behavioural AND security differences tests can't catch."),
            ("migration-critic",
             "Read-only, advisory AI reviewing the migrated code's idiomaticity AND static performance — flags legacy patterns and N+1/bundle/blocking regressions."),
        ],
        col_ratios=[0.22, 0.78],
        header_size=12, body_size=10,
        mono_col0=True,
        row_height=0.215,
    )

    add_text(slide, 0.40, 3.15, 12.5, 0.32,
             "2 Automation Hooks",
             size=15, color=COBALT, font=FONT_HEADER)
    add_table(
        slide, 0.40, 3.48, 12.533, 1.10,
        ["Hook", "What It Does"],
        [
            ("hooks.json",
             "Registers the heartbeat to fire automatically after every file save during migration work."),
            ("heartbeat.mjs",
             "Stamps each developer's in-progress unit in real time so the team can detect stalled work."),
        ],
        col_ratios=[0.22, 0.78],
        header_size=12, body_size=11,
        mono_col0=True,
        row_height=0.31,
    )

    add_text(slide, 0.40, 4.50, 12.5, 0.30,
             "9 Templates  ·  31 Framework Files  ·  4 Workflow Scripts",
             size=14, color=COBALT, font=FONT_HEADER)
    add_table(
        slide, 0.40, 4.82, 12.533, 2.45,
        ["Artifact", "What It Does"],
        [
            ("migration.md",                          "The team-editable configuration: target stack, strategy, auth, cross-cutting concerns, acceptance criteria."),
            ("migration-interview.json",              "Declarative question catalog driving /analyze's interactive setup interview."),
            ("frameworks/<name>.md  ×31",             "One file per supported framework — detection signals, scaffold/integration recipes, test framework, auth notes."),
            ("workflows/analyze-discovery.js",        "Workflow-tool script — exhaustive loop-until-dry entry-point discovery driving /analyze."),
            ("workflows/foundation-establish.js",     "Workflow-tool script — fans out cross-cutting-migrator per concern, in parallel, for /foundation."),
            ("state.schema.json  +  unit.schema.json", "Schemas for the workflow ledger and per-unit state — one JSON file per unit, conflict-free in git."),
            ("plan.md  +  report.md  +  notes-template.md", "Templates for the auto-generated plan, stakeholder reports, and per-unit design notes."),
            ("verify.config.json",                    "Per-stack lint, type-check, test, and dynamic (E2E + API-replay) commands used by /verify."),
            ("permanent-gotchas/fastapi/pyproject.toml", "Reference asset that bypasses a known FastAPI install gotcha automatically."),
        ],
        col_ratios=[0.34, 0.66],
        header_size=12, body_size=9,
        mono_col0=True,
        row_height=0.235,
    )


# ══════════════════════════════════════════════════════════════════════════════
# Slide 13 — Closing / Next Step
# ══════════════════════════════════════════════════════════════════════════════

def slide_13_closing(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 13.333, 7.5, fill=NAVY)
    add_rect(slide, 0, 0, 13.333, 0.07, fill=TEAL)
    add_rect(slide, 0, 7.43, 13.333, 0.07, fill=TEAL)
    add_rect(slide, 0, 0, 0.14, 7.5, fill=TEAL)

    chips = [
        "AI translates · You approve",
        "20+ developers · Zero conflicts",
        "Full visibility · Every sprint",
    ]
    chip_w = 3.6
    chip_gap = 0.30
    total = chip_w * 3 + chip_gap * 2
    chip_start = (13.333 - total) / 2
    for i, text in enumerate(chips):
        x = chip_start + i * (chip_w + chip_gap)
        add_rect(slide, x, 1.20, chip_w, 0.55, fill=COBALT)
        add_text(slide, x + 0.05, 1.22, chip_w - 0.10, 0.51,
                 text, size=13, color=WHITE, font=FONT_HEADER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, 0.7, 2.20, 11.933, 2.0,
             "web-modernize turns a multi-year legacy rewrite\n"
             "into a measurable, sprint-by-sprint delivery\n"
             "with safety nets at every step.",
             size=26, color=WHITE, font=FONT_HEADER,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(slide, 1.0, 4.55, 11.333, 1.50, fill=COBALT)
    add_text(slide, 1.0, 4.65, 11.333, 0.45,
             "Next Step", size=18, color=WHITE, font=FONT_HEADER,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 1.2, 5.15, 11.0, 0.85,
             "Pilot the plugin on one legacy application this sprint.\n"
             "Run  /web-modernize:analyze  and see your full migration scope in minutes.",
             size=15, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, 0.5, 6.30, 12.333, 0.35,
             "Balaji Harikrishnan   ·   balaji.harikrishnan@cognizant.com",
             size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 6.70, 12.333, 0.32,
             "web-modernize v0.18.0   ·   github.com/balaji-hari/web-mordernize",
             size=11, italic=True, color=SLATE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Build
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    print("Building 13-slide leadership deck (v0.18.0)...")
    slide_01_title(prs);                    print("  [1/13]  Title")
    slide_02_intro(prs);                    print("  [2/13]  What is web-modernize")
    slide_09_extensibility(prs);            print("  [3/13]  Natural language + pluggable framework library")
    slide_07_advantages(prs);               print("  [4/13]  What the plugin adds to Claude Code")
    slide_03_need(prs);                     print("  [5/13]  Why we need it")
    slide_08_risk_safety(prs);              print("  [6/13]  Risk, safety, and reversibility")
    slide_10_planning(prs);                 print("  [7/13]  Built for sprint and PI planning")
    slide_04_analyze(prs);                  print("  [8/13]  From legacy codebase to migration backlog")
    slide_05_execution(prs);                print("  [9/13]  From backlog to migrated units")
    slide_06_verification_tiers(prs);       print("  [10/13] Three verification tiers")
    slide_11_commands(prs);                 print("  [11/13] 19 skills (slash commands)")
    slide_12_inventory(prs);                print("  [12/13] Agents, hooks, templates, framework library")
    slide_13_closing(prs);                  print("  [13/13] Closing / next step")

    out = r"C:\1\web-mordernize\docs\decks\web-modernize-presentation.pptx"
    prs.save(out)
    print(f"\nSaved: {out}")


if __name__ == "__main__":
    build()
